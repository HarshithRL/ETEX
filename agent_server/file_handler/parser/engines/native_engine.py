"""Native fallback parsers: PyMuPDF, python-docx, openpyxl, python-pptx."""

from __future__ import annotations

import re
from pathlib import Path

from agent_server.file_handler.parser.types import ParseError, SupportedFormat

_SUPPORTED = frozenset(SupportedFormat)


def _escape_cell(value: object) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\n", " ").replace("|", "\\|").strip()
    return text


def _md_table(rows: list[list[object]]) -> str:
    """Build a GFM table; first row is treated as header when present."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    if width == 0:
        return ""
    normalized: list[list[str]] = []
    for row in rows:
        cells = [_escape_cell(c) for c in row]
        if len(cells) < width:
            cells.extend([""] * (width - len(cells)))
        normalized.append(cells[:width])

    # Drop trailing fully-empty rows
    while normalized and all(c == "" for c in normalized[-1]):
        normalized.pop()
    if not normalized:
        return ""

    header = normalized[0]
    body = normalized[1:] if len(normalized) > 1 else []
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _trim_sheet_matrix(matrix: list[list[object]]) -> list[list[object]]:
    if not matrix:
        return []
    # Drop trailing empty rows
    while matrix and all(
        c is None or (isinstance(c, str) and not c.strip()) for c in matrix[-1]
    ):
        matrix.pop()
    if not matrix:
        return []
    # Drop trailing empty columns
    width = max(len(r) for r in matrix)
    while width > 0:
        col = width - 1
        if all(
            col >= len(r)
            or r[col] is None
            or (isinstance(r[col], str) and not str(r[col]).strip())
            for r in matrix
        ):
            width -= 1
        else:
            break
    return [list(r[:width]) for r in matrix]


class NativeEngine:
    """Lightweight parsers that do not need Docling model weights."""

    name = "native"

    def supports(self, fmt: SupportedFormat) -> bool:
        return fmt in _SUPPORTED

    def parse(self, path: Path, fmt: SupportedFormat) -> str:
        if not self.supports(fmt):
            raise ParseError(f"Unsupported format: {fmt.value}", path=path)
        try:
            if fmt is SupportedFormat.PDF:
                body = self._parse_pdf(path)
            elif fmt is SupportedFormat.DOCX:
                body = self._parse_docx(path)
            elif fmt is SupportedFormat.XLSX:
                body = self._parse_xlsx(path)
            elif fmt is SupportedFormat.PPTX:
                body = self._parse_pptx(path)
            else:  # pragma: no cover
                raise ParseError(f"Unsupported format: {fmt.value}", path=path)
        except ParseError:
            raise
        except Exception as exc:  # noqa: BLE001
            raise ParseError(
                f"Native parse failed: {exc}",
                path=path,
                cause=exc,
            ) from exc

        if not body.strip():
            raise ParseError("Native parser produced empty markdown", path=path)
        return body.strip() + "\n"

    def _parse_pdf(self, path: Path) -> str:
        import pymupdf

        doc = pymupdf.open(path)
        parts: list[str] = []
        try:
            for i, page in enumerate(doc, start=1):
                text = page.get_text("text") or ""
                text = text.strip()
                if not text:
                    continue
                parts.append(f"## Page {i}\n\n{text}")
        finally:
            doc.close()
        return "\n\n---\n\n".join(parts)

    def _parse_docx(self, path: Path) -> str:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = Document(path)
        parts: list[str] = []

        def heading_level(style_name: str | None) -> int | None:
            if not style_name:
                return None
            match = re.match(r"Heading\s+(\d+)", style_name, re.IGNORECASE)
            if match:
                return max(1, min(6, int(match.group(1))))
            if style_name.lower() in {"title", "subtitle"}:
                return 1
            return None

        # Walk body in document order (paragraphs + tables)
        body = document.element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                para = Paragraph(child, document)
                text = para.text.strip()
                if not text:
                    continue
                level = heading_level(para.style.name if para.style else None)
                if level is not None:
                    parts.append("#" * level + " " + text)
                else:
                    parts.append(text)
            elif child.tag == qn("w:tbl"):
                table = Table(child, document)
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                md = _md_table(rows)
                if md:
                    parts.append(md)

        return "\n\n".join(parts)

    def _parse_xlsx(self, path: Path) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(path, data_only=True, read_only=True)
        parts: list[str] = []
        try:
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                matrix: list[list[object]] = []
                for row in sheet.iter_rows(values_only=True):
                    matrix.append(list(row) if row else [])
                matrix = _trim_sheet_matrix(matrix)
                parts.append(f"## Sheet: {sheet_name}")
                if not matrix:
                    parts.append("_(empty)_")
                    continue
                md = _md_table(matrix)
                parts.append(md if md else "_(empty)_")
        finally:
            workbook.close()
        return "\n\n".join(parts)

    def _parse_pptx(self, path: Path) -> str:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(path)
        parts: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = [f"## Slide {index}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts = [
                        p.text.strip()
                        for p in shape.text_frame.paragraphs
                        if p.text and p.text.strip()
                    ]
                    if texts:
                        slide_parts.append("\n".join(texts))
                if shape.has_table:
                    table = shape.table
                    rows = [
                        [cell.text for cell in row.cells] for row in table.rows
                    ]
                    md = _md_table(rows)
                    if md:
                        slide_parts.append(md)
                # Nested group shapes: best-effort text only
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for nested in shape.shapes:
                        if nested.has_text_frame:
                            nested_texts = [
                                p.text.strip()
                                for p in nested.text_frame.paragraphs
                                if p.text and p.text.strip()
                            ]
                            if nested_texts:
                                slide_parts.append("\n".join(nested_texts))
            if len(slide_parts) > 1:
                parts.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(parts)
