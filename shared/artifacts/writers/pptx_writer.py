from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.table import Table
from pptx.util import Inches

from ..exceptions import ArtifactPatchError, ArtifactWriteError
from ..models import ArtifactDocument, ArtifactType, BlockType
from ..patch import LocationTracker, require_table, require_text
from ..spec import ArtifactSpec, PatchOp, SpecBlock, SpecSlide, list_items


class PptxWriter:
    kind = ArtifactType.PPT

    def write(self, spec: ArtifactSpec, dest: Path) -> None:
        try:
            presentation = Presentation()
            if spec.title:
                presentation.core_properties.title = spec.title
            if spec.author:
                presentation.core_properties.author = spec.author
            slides = list(spec.slides or [])
            if not slides:
                slides = [SpecSlide(blocks=[SpecBlock(type=BlockType.HEADING, text=spec.title or "Untitled", level=1)])]
            for index, slide_spec in enumerate(slides):
                _add_slide(presentation, slide_spec, is_title=index == 0)
            dest.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(str(dest))
        except ArtifactWriteError:
            raise
        except Exception as exc:
            raise ArtifactWriteError(f"Failed to write PowerPoint file {dest.name}: {exc}") from exc

    def apply_ops(
        self,
        path: Path,
        document: ArtifactDocument,
        ops: list[PatchOp],
        dest: Path,
    ) -> None:
        presentation = Presentation(str(path))
        tracker = LocationTracker(document.blocks)
        for op in ops:
            _apply_pptx_op(presentation, tracker, op)
        dest.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(dest))


def _add_slide(presentation: Presentation, slide_spec: SpecSlide, *, is_title: bool) -> None:
    heading = next((block for block in slide_spec.blocks if block.type == BlockType.HEADING), None)
    rest = [block for block in slide_spec.blocks if block is not heading]
    title_text = (heading.text if heading else "") or ""
    texts: list[str] = []
    tables: list[list[list[str]]] = []
    for block in rest:
        if block.type == BlockType.TABLE and block.table:
            tables.append(block.table)
        elif block.type == BlockType.LIST:
            texts.extend(list_items(block))
        elif block.type == BlockType.NOTE:
            continue
        elif block.text.strip():
            texts.extend(line.strip() for line in block.text.splitlines() if line.strip())

    if is_title and not tables and not rest:
        layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title_text
        if len(slide.placeholders) > 1 and texts:
            slide.placeholders[1].text = "\n".join(texts)
        _set_notes(slide, slide_spec.notes)
        return

    layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title_text
    body = _body_placeholder(slide)
    if body is not None and texts:
        frame = body.text_frame
        frame.clear()
        frame.paragraphs[0].text = texts[0]
        for line in texts[1:]:
            paragraph = frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
    elif texts:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(2.2))
        box.text_frame.text = "\n".join(texts)
    table_top = Inches(3.4 if texts else 1.8)
    for table in tables:
        _add_table(slide, table, top=table_top)
        table_top += Inches(0.4) + Inches(0.32 * max(len(table), 1))
    _set_notes(slide, slide_spec.notes)


def _body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape
    if len(slide.placeholders) > 1:
        return slide.placeholders[1]
    return None


def _add_table(slide, rows: list[list[str]], *, top=None) -> None:
    n_rows = len(rows)
    n_cols = max(len(row) for row in rows)
    width = Inches(9.0)
    height = Inches(max(0.8, 0.32 * n_rows))
    left = Inches(0.5)
    if top is None:
        top = Inches(3.4)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    _fill_ppt_table(table, rows)


def _fill_ppt_table(table: Table, rows: list[list[str]]) -> None:
    for r_i, row in enumerate(rows):
        for c_i, value in enumerate(row):
            if r_i < len(table.rows) and c_i < len(table.columns):
                table.cell(r_i, c_i).text = value


def _set_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def _apply_pptx_op(presentation: Presentation, tracker: LocationTracker, op: PatchOp) -> None:
    if op.op == "set_title":
        title = op.title or op.text
        if title is None:
            raise ArtifactPatchError("set_title requires 'title' or 'text'.")
        presentation.core_properties.title = title
        if presentation.slides:
            first = presentation.slides[0]
            if first.shapes.title is not None and (op.block_id is None):
                first.shapes.title.text = title
        return
    if op.op == "set_cell":
        raise ArtifactPatchError("set_cell is only valid for excel artifacts.")
    if op.op in {"replace_text", "replace_shape_text"}:
        _replace_shape_text(presentation, tracker, op)
        return
    if op.op == "replace_table":
        _replace_shape_table(presentation, tracker, op)
        return
    if op.op == "insert_slide":
        if op.slide is None:
            raise ArtifactPatchError("insert_slide requires a 'slide' payload.")
        _add_slide(presentation, op.slide, is_title=False)
        return
    if op.op == "insert_block":
        _insert_block(presentation, tracker, op)
        return
    if op.op == "delete_block":
        _delete_block(presentation, tracker, op)
        return
    raise ArtifactPatchError(f"Unknown patch op {op.op!r}.")


def _replace_shape_text(presentation: Presentation, tracker: LocationTracker, op: PatchOp) -> None:
    block = tracker.find(op)
    slide_no = block.location.slide
    if slide_no is None:
        raise ArtifactPatchError(f"Block {block.block_id} has no slide location.")
    slide = presentation.slides[slide_no - 1]
    if block.type == BlockType.NOTE or not block.location.shape_id:
        slide.notes_slide.notes_text_frame.text = require_text(op)
        return
    shape = _shape_by_id(slide, block.location.shape_id)
    if shape is None or not getattr(shape, "has_text_frame", False):
        raise ArtifactPatchError(f"No text shape with id {block.location.shape_id!r} on slide {slide_no}.")
    shape.text_frame.text = require_text(op)


def _replace_shape_table(presentation: Presentation, tracker: LocationTracker, op: PatchOp) -> None:
    block = tracker.find(op)
    slide_no = block.location.slide
    if slide_no is None or not block.location.shape_id:
        raise ArtifactPatchError(f"Block {block.block_id} has no slide/shape_id for table patch.")
    slide = presentation.slides[slide_no - 1]
    shape = _shape_by_id(slide, block.location.shape_id)
    rows = require_table(op)
    if shape is None or not getattr(shape, "has_table", False):
        raise ArtifactPatchError(f"No table shape with id {block.location.shape_id!r} on slide {slide_no}.")
    n_rows = len(rows)
    n_cols = max((len(row) for row in rows), default=0)
    table = shape.table
    if len(table.rows) == n_rows and len(table.columns) == n_cols:
        _fill_ppt_table(table, rows)
        return
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    _remove_shape(shape)
    new_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    _fill_ppt_table(new_shape.table, rows)


def _insert_block(presentation: Presentation, tracker: LocationTracker, op: PatchOp) -> None:
    if op.block is None:
        raise ArtifactPatchError("insert_block requires a 'block' payload.")
    if not op.after_block_id and not op.block_id and op.location is None:
        _add_slide(
            presentation,
            SpecSlide(blocks=[op.block]),
            is_title=False,
        )
        return
    anchor = tracker.find_anchor(op)
    slide_no = anchor.location.slide or 1
    slide = presentation.slides[slide_no - 1]
    block = op.block
    if block.type == BlockType.TABLE and block.table:
        _add_table(slide, block.table, top=Inches(5.0))
        return
    box = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(8.8), Inches(1.2))
    text = "\n".join(list_items(block) or [block.text])
    box.text_frame.text = text


def _delete_block(presentation: Presentation, tracker: LocationTracker, op: PatchOp) -> None:
    block = tracker.find(op)
    slide_no = block.location.slide
    if slide_no is None:
        raise ArtifactPatchError(f"Block {block.block_id} has no slide location.")
    slide = presentation.slides[slide_no - 1]
    if block.type == BlockType.NOTE or not block.location.shape_id:
        slide.notes_slide.notes_text_frame.text = ""
        tracker.remove(block)
        return
    shape = _shape_by_id(slide, block.location.shape_id)
    if shape is None:
        raise ArtifactPatchError(f"No shape with id {block.location.shape_id!r} on slide {slide_no}.")
    _remove_shape(shape)
    tracker.remove(block)


def _shape_by_id(slide, shape_id: str):
    for shape in slide.shapes:
        if str(getattr(shape, "shape_id", "") or "") == str(shape_id):
            return shape
    return None


def _remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)
