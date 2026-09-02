from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..blocks import make_content_block
from ..exceptions import ArtifactParseError, CorruptArtifact
from ..ids import file_artifact_id
from ..models import (
    ArtifactDocument,
    ArtifactMetadata,
    ArtifactType,
    BlockType,
    ContentBlock,
    MIME_BY_TYPE,
    ParseOptions,
    SourceLocation,
)
from .pdf_tables import table_to_text


class PptxParser:
    kind = ArtifactType.PPT

    def parse(self, path: Path, options: ParseOptions) -> ArtifactDocument:
        try:
            presentation = Presentation(str(path))
        except PackageNotFoundError as exc:
            raise CorruptArtifact(f"Unreadable PowerPoint file: {path.name}") from exc
        except Exception as exc:
            raise CorruptArtifact(f"Failed to open PowerPoint file: {path.name}") from exc

        core = presentation.core_properties
        slides = list(presentation.slides)
        if options.max_pages is not None:
            slides = slides[: options.max_pages]

        blocks: list[ContentBlock] = []
        warnings: list[str] = []
        index = 0

        try:
            for slide_i, slide in enumerate(slides, start=1):
                for shape in slide.shapes:
                    shape_id = str(getattr(shape, "shape_id", "") or "")
                    if getattr(shape, "has_text_frame", False):
                        texts = [
                            para.text.strip()
                            for para in shape.text_frame.paragraphs
                            if para.text and para.text.strip()
                        ]
                        if texts:
                            index += 1
                            blocks.append(
                                make_content_block(
                                    seq_id=f"ppt-{index:04d}",
                                    block_type=BlockType.TEXT,
                                    text="\n".join(texts),
                                    location=SourceLocation(slide=slide_i, shape_id=shape_id or None),
                                )
                            )
                    if options.include_tables and getattr(shape, "has_table", False):
                        rows = [
                            [cell.text.strip() for cell in row.cells]
                            for row in shape.table.rows
                        ]
                        if any(cell for row in rows for cell in row):
                            index += 1
                            blocks.append(
                                make_content_block(
                                    seq_id=f"ppt-{index:04d}",
                                    block_type=BlockType.TABLE,
                                    text=table_to_text(rows),
                                    table=rows,
                                    location=SourceLocation(slide=slide_i, shape_id=shape_id or None),
                                )
                            )
                    if options.include_images and shape.shape_type is not None:
                        if int(shape.shape_type) == 13:
                            index += 1
                            blocks.append(
                                make_content_block(
                                    seq_id=f"ppt-{index:04d}",
                                    block_type=BlockType.IMAGE,
                                    text=f"[image on slide {slide_i}]",
                                    location=SourceLocation(slide=slide_i, shape_id=shape_id or None),
                                )
                            )

                notes_text = _slide_notes(slide)
                if notes_text:
                    index += 1
                    blocks.append(
                        make_content_block(
                            seq_id=f"ppt-{index:04d}",
                            block_type=BlockType.NOTE,
                            text=notes_text,
                            location=SourceLocation(slide=slide_i),
                        )
                    )
        except Exception as exc:
            raise ArtifactParseError(f"Failed to extract PowerPoint {path.name}: {exc}") from exc

        if not blocks:
            warnings.append("PowerPoint contained no extractable text, tables, or notes.")

        return ArtifactDocument(
            source=str(path),
            artifact_type=ArtifactType.PPT,
            artifact_id=file_artifact_id(path),
            metadata=ArtifactMetadata(
                filename=path.name,
                artifact_type=ArtifactType.PPT,
                mime_type=MIME_BY_TYPE[ArtifactType.PPT],
                title=_meta_str(core.title),
                author=_meta_str(core.author),
                created=_meta_str(core.created),
                modified=_meta_str(core.modified),
                slide_count=len(presentation.slides),
            ),
            blocks=blocks,
            warnings=warnings,
        )


def _slide_notes(slide) -> str:
    notes_slide = getattr(slide, "has_notes_slide", False) and slide.notes_slide
    if not notes_slide or notes_slide.notes_text_frame is None:
        return ""
    return "\n".join(
        para.text.strip()
        for para in notes_slide.notes_text_frame.paragraphs
        if para.text and para.text.strip()
    )


def _meta_str(value: object) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    return text or None
