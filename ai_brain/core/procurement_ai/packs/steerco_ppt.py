"""SteerCo deck generated only from named Excel fields."""

from __future__ import annotations

from typing import Any

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ai_brain.core.procurement_ai.packs import store
from ai_brain.core.procurement_ai.packs.compare_xlsx import existing_xlsx_path

ORANGE = RGBColor(0xFF, 0x69, 0x00)
DARK = RGBColor(0x20, 0x20, 0x20)
MUTED = RGBColor(0x66, 0x66, 0x66)


def build_steerco_ppt(project_id: str, *, thread_id: str = "") -> dict[str, Any]:
    xlsx_path = existing_xlsx_path(project_id)
    if xlsx_path is None:
        store.write_status(project_id, "ppt", status="blocked", error="Comparison Excel is not ready", thread_id=thread_id or None)
        return {"status": "blocked", "error": "Comparison Excel is not ready"}
    store.write_status(project_id, "ppt", status="running", thread_id=thread_id or None)
    fields = _named_fields(xlsx_path)
    vendors = _vendor_rows(xlsx_path)
    filename = "steerco_pack.pptx"
    path = store.pack_file(project_id, filename)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _title_slide(prs, fields.get("FIELD_PROJECT_NAME") or "Vendor comparison", fields.get("FIELD_PROCESS_LABEL") or "", fields.get("FIELD_OWNER_ENTITY") or "", fields.get("FIELD_PROJECT_CODE") or "")
    _bullets_slide(prs, "Process and ownership", [
        f"Process: {fields.get('FIELD_PROCESS_LABEL') or 'unset'}",
        f"Owner: {fields.get('FIELD_OWNER_ENTITY') or 'unset'}",
        f"Knowledge coverage: {fields.get('FIELD_KNOWLEDGE_PCT') or '—'}%",
        f"Weights: {fields.get('FIELD_WEIGHTS_LOCKED') or 'unlocked'}",
    ])
    _vendor_slide(prs, vendors)
    _bullets_slide(prs, "Decision state", [
        fields.get("FIELD_DECISION_SUMMARY") or "No summary field on Cover!B12",
        f"Award: {fields.get('FIELD_AWARD') or 'Not awarded'}",
        "Agents draft. Humans award.",
    ])
    prs.save(path)
    href = store.flask_href(project_id, "ppt")
    store.write_status(project_id, "ppt", status="ready", href=href, filename=filename, thread_id=thread_id or None, error=None)
    return {"status": "ready", "href": href, "filename": filename, "path": str(path)}


def _named_fields(xlsx_path) -> dict[str, str]:
    cover = load_workbook(xlsx_path, data_only=True)["Cover"]
    mapping = {
        "FIELD_PROCESS_TYPE": cover["B5"].value,
        "FIELD_PROCESS_LABEL": cover["B6"].value,
        "FIELD_OWNER_ENTITY": cover["B7"].value,
        "FIELD_PROJECT_CODE": cover["B8"].value,
        "FIELD_PROJECT_NAME": cover["B9"].value,
        "FIELD_KNOWLEDGE_PCT": cover["B10"].value,
        "FIELD_KB_STATUS": cover["B11"].value,
        "FIELD_DECISION_SUMMARY": cover["B12"].value,
        "FIELD_WEIGHTS_LOCKED": cover["B13"].value,
        "FIELD_AWARD": cover["B14"].value,
    }
    return {key: "" if value is None else str(value) for key, value in mapping.items()}


def _vendor_rows(xlsx_path) -> list[tuple[str, str]]:
    sheet = load_workbook(xlsx_path, data_only=True)["Vendors"]
    rows = []
    for row in sheet.iter_rows(min_row=2, values_only=True):
        name = row[0]
        headline = row[4] if len(row) > 4 else ""
        if name:
            rows.append((str(name), str(headline or "")))
    return rows


def _title_slide(prs, title: str, process: str, owner: str, code: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, title, 0.6, 2.4, 12, 1.2, size=36, bold=True, color=DARK)
    _box(slide, "SteerCo pack — generated from comparison Excel", 0.6, 3.6, 12, 0.4, size=16, color=ORANGE)
    _box(slide, " · ".join(part for part in (code, process, owner) if part), 0.6, 4.2, 12, 0.4, size=14, color=MUTED)


def _vendor_slide(prs, vendors: list[tuple[str, str]]) -> None:
    lines = [f"{name} — {headline}" if headline else name for name, headline in vendors] or ["No vendor columns in Vendors!A"]
    _bullets_slide(prs, "Vendor glance", lines[:8])


def _bullets_slide(prs, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, title, 0.6, 0.4, 12, 0.6, size=24, bold=True, color=DARK)
    text = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(5.6))
    tf = text.text_frame
    tf.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = DARK
        paragraph.space_after = Pt(10)


def _box(slide, text: str, left: float, top: float, width: float, height: float, *, size: int, bold: bool = False, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
